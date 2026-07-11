"""
ai_assistant app - Module 8: AI Study Suggestions, Chatbot & Quiz views.

Endpoints:
    GET  /api/ai/suggestions/                 - list saved suggestions
    POST /api/ai/suggestions/generate/        - generate new daily suggestions
    PATCH /api/ai/suggestions/<id>/           - mark a suggestion done
    GET  /api/ai/chatbot/history/             - chatbot conversation history
    POST /api/ai/chatbot/ask/                 - send a message, get AI reply
    GET  /api/ai/quizzes/                     - list past quizzes
    POST /api/ai/quizzes/generate/            - generate a new AI quiz
    GET  /api/ai/quizzes/<id>/                - view a quiz (questions, no answers if unattempted)
    POST /api/ai/quizzes/<id>/submit/         - submit answers, get scored result
"""
from django.shortcuts import get_object_or_404
from rest_framework import generics, permissions, status, viewsets
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework.parsers import MultiPartParser, FormParser
import fitz
from pptx import Presentation

from profiles.models import StudentProfile, Subject
from .models import StudySuggestion, ChatbotMessage, Quiz, QuizQuestion
from .serializers import (
    StudySuggestionSerializer, ChatbotMessageSerializer, ChatbotAskSerializer,
    QuizSerializer, QuizGenerateSerializer, QuizSubmitSerializer,
    QuizQuestionResultSerializer,
)
from . import services


class StudySuggestionViewSet(viewsets.ModelViewSet):
    """List, update (mark done), delete saved suggestions."""
    serializer_class = StudySuggestionSerializer
    permission_classes = [permissions.IsAuthenticated]
    http_method_names = ['get', 'patch', 'delete', 'head', 'options']

    def get_queryset(self):
        return StudySuggestion.objects.filter(student=self.request.user)


class GenerateSuggestionsView(APIView):
    """POST /api/ai/suggestions/generate/ - Module 8 AI suggestion generation."""
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request):
        me = request.user
        profile, _ = StudentProfile.objects.get_or_create(user=me)
        weak_subjects = list(profile.weak_subjects.all())

        texts = services.generate_daily_suggestions(profile, weak_subjects, count=3)
        suggestions = [
            StudySuggestion.objects.create(student=me, text=text)
            for text in texts
        ]
        serializer = StudySuggestionSerializer(suggestions, many=True)
        return Response(serializer.data, status=status.HTTP_201_CREATED)


class ChatbotHistoryView(generics.ListAPIView):
    """GET /api/ai/chatbot/history/"""
    serializer_class = ChatbotMessageSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        return ChatbotMessage.objects.filter(student=self.request.user)


class ChatbotAskView(APIView):
    """POST /api/ai/chatbot/ask/  body: {"message": "..."}"""
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request):
        serializer = ChatbotAskSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        user_message = serializer.validated_data['message']
        me = request.user

        ChatbotMessage.objects.create(student=me, role='user', content=user_message)

        history = list(
            ChatbotMessage.objects.filter(student=me).order_by('created_at')
            .values('role', 'content')[:50]
        )
        reply_text = services.chatbot_reply(me, history, user_message)

        reply_msg = ChatbotMessage.objects.create(student=me, role='assistant', content=reply_text)
        return Response(ChatbotMessageSerializer(reply_msg).data, status=status.HTTP_201_CREATED)


class QuizListView(generics.ListAPIView):
    """GET /api/ai/quizzes/"""
    serializer_class = QuizSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        return Quiz.objects.filter(student=self.request.user)


class QuizDetailView(generics.RetrieveAPIView):
    """GET /api/ai/quizzes/<id>/"""
    serializer_class = QuizSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        return Quiz.objects.filter(student=self.request.user)


class GenerateQuizView(APIView):
    """POST /api/ai/quizzes/generate/ - Module 8 AI quiz generation after a study session."""
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request):
        serializer = QuizGenerateSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        data = serializer.validated_data

        subject = None
        if data.get('subject'):
            subject = get_object_or_404(Subject, pk=data['subject'])

        quiz = Quiz.objects.create(student=request.user, subject=subject, topic=data['topic'])
        questions_data = services.generate_quiz_questions(data['topic'], data['num_questions'])

        for q in questions_data:
            QuizQuestion.objects.create(
                quiz=quiz,
                question_text=q['question_text'],
                option_a=q['option_a'], option_b=q['option_b'],
                option_c=q['option_c'], option_d=q['option_d'],
                correct_option=q['correct_option'],
            )

        return Response(QuizSerializer(quiz).data, status=status.HTTP_201_CREATED)


class SubmitQuizView(APIView):
    """POST /api/ai/quizzes/<id>/submit/ body: {"answers": {"<question_id>": "a", ...}}"""
    permission_classes = [permissions.IsAuthenticated]

    def post(self, request, pk):
        quiz = get_object_or_404(Quiz, pk=pk, student=request.user)
        serializer = QuizSubmitSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        answers = serializer.validated_data['answers']

        correct_count = 0
        questions = quiz.questions.all()
        for question in questions:
            answer = answers.get(str(question.id))
            if answer:
                question.student_answer = answer
                question.save()
            if question.is_correct():
                correct_count += 1

        quiz.score = round((correct_count / questions.count()) * 100) if questions.count() else 0
        quiz.save()

        return Response({
            'quiz_id': quiz.id,
            'score': quiz.score,
            'correct_count': correct_count,
            'total_questions': questions.count(),
            'questions': QuizQuestionResultSerializer(questions, many=True).data,
        })

class GenerateFileQuizView(APIView):
    """POST /api/ai/quizzes/generate-from-file/ - AI quiz generation from a PDF/PPT file."""
    permission_classes = [permissions.IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request):
        file_obj = request.FILES.get('file')
        if not file_obj:
            return Response({"error": "No file provided."}, status=status.HTTP_400_BAD_REQUEST)
        
        num_questions = int(request.data.get('num_questions', 5))
        topic = request.data.get('topic', 'Uploaded Document')
        
        extracted_text = ""
        filename = file_obj.name.lower()
        
        try:
            if filename.endswith('.pdf'):
                doc = fitz.open(stream=file_obj.read(), filetype="pdf")
                for page in doc:
                    extracted_text += page.get_text() + "\n"
            elif filename.endswith('.ppt') or filename.endswith('.pptx'):
                prs = Presentation(file_obj)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            else:
                return Response({"error": "Unsupported file format. Please upload PDF or PPT."}, status=status.HTTP_400_BAD_REQUEST)
        except Exception as e:
            return Response({"error": f"Error parsing file: {str(e)}"}, status=status.HTTP_400_BAD_REQUEST)
        
        if not extracted_text.strip():
            return Response({"error": "No text could be extracted from the file."}, status=status.HTTP_400_BAD_REQUEST)
            
        # truncate text if too long (AI context limit)
        extracted_text = extracted_text[:10000]
        
        quiz = Quiz.objects.create(student=request.user, topic=topic)
        prompt_topic = f"the following text extracted from a document: \\n\\n{extracted_text}"
        questions_data = services.generate_quiz_questions(prompt_topic, num_questions)

        for q in questions_data:
            QuizQuestion.objects.create(
                quiz=quiz,
                question_text=q['question_text'],
                option_a=q['option_a'], option_b=q['option_b'],
                option_c=q['option_c'], option_d=q['option_d'],
                correct_option=q['correct_option'],
            )

        return Response(QuizSerializer(quiz).data, status=status.HTTP_201_CREATED)

class GenerateFlashcardsView(APIView):
    """POST /api/ai/flashcards/generate/ - AI flashcard generation from topic or file."""
    permission_classes = [permissions.IsAuthenticated]
    # Allow multipart/form-data for files or regular json
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request):
        topic = request.data.get('topic', '')
        num_cards = int(request.data.get('num_cards', 10))
        
        file_obj = request.FILES.get('file')
        extracted_text = ""
        
        if file_obj:
            filename = file_obj.name.lower()
            topic = request.data.get('topic', 'Uploaded Document')
            try:
                if filename.endswith('.pdf'):
                    doc = fitz.open(stream=file_obj.read(), filetype="pdf")
                    for page in doc:
                        extracted_text += page.get_text() + "\n"
                elif filename.endswith('.ppt') or filename.endswith('.pptx'):
                    prs = Presentation(file_obj)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text += shape.text + "\n"
                else:
                    return Response({"error": "Unsupported file format. Please upload PDF or PPT."}, status=status.HTTP_400_BAD_REQUEST)
            except Exception as e:
                return Response({"error": f"Error parsing file: {str(e)}"}, status=status.HTTP_400_BAD_REQUEST)
                
            extracted_text = extracted_text[:10000]
            prompt_topic = f"the following text extracted from a document: \\n\\n{extracted_text}"
        else:
            prompt_topic = topic
            if not topic:
                return Response({"error": "Topic or file is required."}, status=status.HTTP_400_BAD_REQUEST)
        
        cards = services.generate_flashcards(prompt_topic, num_cards)
        return Response(cards, status=status.HTTP_200_OK)
